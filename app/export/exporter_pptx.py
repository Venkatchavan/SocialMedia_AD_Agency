"""PPTX exporter — renders brief data as a PowerPoint presentation.

Uses python-pptx. Slides: Title, Overview, Hooks, Angles, Scripts, Brand Voice.
"""

from __future__ import annotations

from typing import Any

from app.export.base_exporter import BaseExporter, ExportResult


class PptxExporter(BaseExporter):
    """Export briefs as PowerPoint (.pptx) files."""

    format_name = "pptx"

    def export(self, brief_data: dict[str, Any], output_path: str) -> ExportResult:
        """Render brief data to a PPTX file."""
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.util import Inches, Pt  # type: ignore[import-untyped]
        except ImportError:
            return ExportResult(
                format="pptx",
                output_path="",
                success=False,
                error="python-pptx not installed. pip install python-pptx",
            )

        path = self._ensure_dir(f"{output_path}.pptx")
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        title = self._safe_get(brief_data, "title", "Creative Brief")
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        subtitle = slide.placeholders[1]
        ws = self._safe_get(brief_data, "workspace_id", "")
        subtitle.text = f"Workspace: {ws}" if ws else "Creative Intelligence OS"

        # Overview slide
        overview = self._safe_get(brief_data, "overview")
        if overview:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Overview"
            slide.placeholders[1].text = str(overview)

        # Hooks slide
        hooks = self._safe_get(brief_data, "hooks", [])
        if hooks:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Top Hooks"
            body = slide.placeholders[1].text_frame
            body.clear()
            for hook in hooks:
                text = hook.get("text", hook) if isinstance(hook, dict) else str(hook)
                p = body.add_paragraph()
                p.text = text
                p.font.size = Pt(18)

        # Angles slide
        angles = self._safe_get(brief_data, "angles", [])
        if angles:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Content Angles"
            body = slide.placeholders[1].text_frame
            body.clear()
            for angle in angles:
                if isinstance(angle, dict):
                    text = f"{angle.get('name', '')}: {angle.get('description', '')}"
                else:
                    text = str(angle)
                p = body.add_paragraph()
                p.text = text

        # Scripts slide(s) — one per script
        scripts = self._safe_get(brief_data, "scripts", [])
        for idx, script in enumerate(scripts, 1):
            if not isinstance(script, dict):
                continue
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = f"Script {idx}"
            body = slide.placeholders[1].text_frame
            body.clear()
            hook_text = script.get("hook", "")
            if hook_text:
                p = body.add_paragraph()
                p.text = f"Hook: {hook_text}"
                p.font.bold = True
            for scene in script.get("scenes", []):
                p = body.add_paragraph()
                p.text = (f"Scene {scene.get('scene_number', '')}: "
                          f"{scene.get('dialogue', '')}")

        # Brand voice slide
        brand_voice = self._safe_get(brief_data, "brand_voice")
        if brand_voice:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = "Brand Voice"
            body = slide.placeholders[1].text_frame
            body.clear()
            if isinstance(brand_voice, dict):
                for k, v in brand_voice.items():
                    p = body.add_paragraph()
                    p.text = f"{k}: {v}"
            else:
                body.add_paragraph().text = str(brand_voice)

        prs.save(str(path))
        file_size = path.stat().st_size

        return ExportResult(
            format="pptx",
            output_path=str(path),
            file_size_bytes=file_size,
        )
